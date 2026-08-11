"""Per-page extracted rich text (<b>/<i>/<p>/<br> only).

Native text layer (PyMuPDF spans / python-pptx runs) when it exists —
preserving bold/italic. For scanned pages with no text layer, falls back to
the stored Docling-OCR elements, so scanned PDFs finally become selectable
and searchable. Everything is HTML-escaped before tagging.
"""

import html
import re

from manabi_core.models import DocElement, Document, DocumentKind, DocumentPage
from sqlalchemy import select
from sqlalchemy.orm import Session

from manabi_server.storage import files

USER_EDITED_MARKER = "<!--user-edited-->"

_BLOCK_RE = re.compile(
    r"<p\b[^>]*>.*?</p>|<table\b[^>]*>.*?</table>|<pre\b[^>]*>.*?</pre>",
    re.DOTALL | re.IGNORECASE,
)
_INNER_P_RE = re.compile(r"(?is)^<p\b[^>]*>(.*)</p>\s*$")


def _block_plain(block: str) -> str:
    return html.unescape(re.sub(r"<[^>]+>", " ", block)).strip()


def _merge_paragraphs(a: str, b: str) -> str:
    """Splice two <p> blocks into one, dehyphenating a split word."""
    ia = (_INNER_P_RE.sub(r"\1", a)).rstrip()
    ib = (_INNER_P_RE.sub(r"\1", b)).lstrip()
    if _block_plain(a).endswith("-"):
        ia = re.sub(r"-\s*$", "", ia)
        return f"<p>{ia}{ib}</p>"
    return f"<p>{ia} {ib}</p>"


def merge_pages_html(page_htmls: list[str | None]) -> str:
    """One continuous document: every page's text_html concatenated, with a
    paragraph that continues onto the next page joined into one flowing block
    (so cross-page breaks don't read as info on 'the wrong page'). Pure function
    over the ordered per-page HTML."""
    from manabi_server.processing.textmerge import should_join

    out: list[tuple[str, str]] = []  # (kind, html)
    for text_html in page_htmls:
        page_html = text_html or ""
        if page_html.startswith(USER_EDITED_MARKER):
            page_html = page_html[len(USER_EDITED_MARKER) :]
        for m in _BLOCK_RE.finditer(page_html):
            block = m.group(0)
            kind = "p" if block[:2].lower() == "<p" else "other"
            if (
                kind == "p"
                and out
                and out[-1][0] == "p"
                and should_join(_block_plain(out[-1][1]), _block_plain(block))
            ):
                out[-1] = ("p", _merge_paragraphs(out[-1][1], block))
            else:
                out.append((kind, block))
    return "".join(b for _, b in out)


BOLD_FLAG = 1 << 4  # PyMuPDF span flag bit for bold
ITALIC_FLAG = 1 << 1
MONO_FLAG = 1 << 3  # PyMuPDF: monospaced font
_MONO_FONT = ("courier", "mono", "consolas", "typewriter")


# Postgres text columns reject NUL; some PDF fonts leak NUL/control chars
_CONTROL_CHARS = re.compile(r"[\x00-\x08\x0b\x0c\x0e-\x1f]")

# Runs of horizontal whitespace (spaces/tabs). PDF and OCR extraction pad
# justified text with several spaces between words and tabs mid-line; those
# read as ragged gaps and make selections mismatch the visible text.
_HSPACE_RUN = re.compile(r"[ \t]+")


def sanitize(text: str) -> str:
    return _CONTROL_CHARS.sub("", text)


def collapse_hspace(text: str) -> str:
    """Collapse runs of spaces/tabs to a single space. Newlines are preserved
    so line structure survives; leading/trailing trimming is the caller's job."""
    return _HSPACE_RUN.sub(" ", text)


# ── Repeated-line (boilerplate) detection ─────────────────────────────────
# Shared with pipeline._strip_boilerplate so the elements/chunks cleanup and
# the native text views agree on what counts as boilerplate.

REPEAT_PAGE_RATIO = 0.4  # line must appear on ≥ this fraction of pages
REPEAT_MAX_LEN = 80  # only short lines (footers, page numbers) qualify
REPEAT_MIN_PAGES = 5  # never strip from tiny documents


def repeat_key(text: str) -> str:
    """Normalization for frequency counting: digits → # so 'Page 3' and
    'Page 57' collapse to one key."""
    return re.sub(r"\d+", "#", text.strip().lower())


def repeated_keys(per_page_keys: dict[int, set[str]]) -> set[str]:
    """Keys of short lines appearing on ≥ REPEAT_PAGE_RATIO of pages."""
    if len(per_page_keys) < REPEAT_MIN_PAGES:
        return set()
    counts: dict[str, int] = {}
    for keys in per_page_keys.values():
        for key in keys:
            counts[key] = counts.get(key, 0) + 1
    total = len(per_page_keys)
    return {k for k, c in counts.items() if c / total >= REPEAT_PAGE_RATIO}


def _is_repeat(raw: str, skip: set[str] | None) -> bool:
    if not skip:
        return False
    raw = raw.strip()
    return bool(raw) and len(raw) <= REPEAT_MAX_LEN and repeat_key(raw) in skip


def _span_html(text: str, bold: bool, italic: bool, mono: bool = False) -> str:
    out = html.escape(sanitize(text))
    if mono:
        out = f"<code>{out}</code>"  # code samples keep their monospace look
    if bold:
        out = f"<b>{out}</b>"
    if italic:
        out = f"<i>{out}</i>"
    return out


INDENT_STEP_PTS = 18  # one indent level ≈ 18pt of left offset
MAX_INDENT = 4


def _pdf_page_line_keys(page) -> set[str]:
    """Normalized keys of this page's short native lines (for boilerplate
    frequency counting across the document)."""
    keys: set[str] = set()
    for block in page.get_text("dict").get("blocks", []):
        if block.get("type") != 0:
            continue
        for line in block.get("lines", []):
            raw = "".join(s.get("text", "") for s in line.get("spans", [])).strip()
            if raw and len(raw) <= REPEAT_MAX_LEN:
                keys.add(repeat_key(raw))
    return keys


def _pdf_page_html(page, skip: set[str] | None = None) -> str | None:
    """Styled HTML from the native text layer; None when the page is scanned.
    Preserves bold/italic and approximates indentation from span x-origins.
    Lines whose normalized text is in `skip` (document-wide boilerplate:
    copyright footers, page numbers) are dropped."""
    data = page.get_text("dict")
    text_blocks = [b for b in data.get("blocks", []) if b.get("type") == 0]

    # Base left margin for indent bucketing
    x_origins = [
        line["bbox"][0]
        for block in text_blocks
        for line in block.get("lines", [])
        if line.get("spans")
    ]
    base_x = min(x_origins) if x_origins else 0.0

    parts: list[str] = []
    for block in text_blocks:
        lines: list[str] = []
        block_x = None
        for line in block.get("lines", []):
            raw_line = "".join(s.get("text", "") for s in line.get("spans", []))
            if _is_repeat(raw_line, skip):
                continue
            spans = []
            prev_trailing_space = True  # start-of-line → trim any leading spaces
            for span in line.get("spans", []):
                text = span.get("text", "")
                if not text.strip():
                    continue
                text = collapse_hspace(text)
                if prev_trailing_space:
                    text = text.lstrip(" ")  # avoid a double space across spans
                prev_trailing_space = text.endswith(" ")
                flags = span.get("flags", 0)
                font = span.get("font", "").lower()
                bold = bool(flags & BOLD_FLAG) or "bold" in font
                italic = bool(flags & ITALIC_FLAG) or "italic" in font
                mono = bool(flags & MONO_FLAG) or any(m in font for m in _MONO_FONT)
                spans.append(_span_html(text, bold, italic, mono))
            if spans:
                if block_x is None:
                    block_x = line["bbox"][0]
                lines.append("".join(spans))
        if lines:
            indent = 0
            if block_x is not None:
                indent = min(MAX_INDENT, int((block_x - base_x) / INDENT_STEP_PTS))
            style = f' style="margin-left:{indent}em"' if indent > 0 else ""
            parts.append(f"<p{style}>{'<br>'.join(lines)}</p>")
    joined = "".join(parts)
    return joined if joined.strip() else None


def _figure_markers(elements: list[DocElement]) -> str:
    """Placeholder lines so readers know visuals exist on the page."""
    markers = []
    for el in elements:
        if el.element_type == "figure":
            markers.append("<p><i>[Figure — see rendered page]</i></p>")
        elif el.element_type == "table" and not (el.text_content or "").strip():
            markers.append("<p><i>[Table — see rendered page]</i></p>")
    return "".join(markers)


def _elements_page_html(elements: list[DocElement]) -> str | None:
    """OCR fallback: build from stored extraction. Headings render bold;
    consecutive paragraph fragments are merged back into real paragraphs
    (OCR splits sentences at line breaks and hyphens)."""
    from manabi_server.processing.textmerge import merge_fragments

    parts: list[str] = []
    paragraph_run: list[str] = []

    def flush_run() -> None:
        for para in merge_fragments(paragraph_run):
            escaped = html.escape(collapse_hspace(para)).replace("\n", "<br>")
            parts.append(f"<p>{escaped}</p>")
        paragraph_run.clear()

    for el in elements:
        if el.element_type == "table" and el.table_json:
            flush_run()
            table_html = _table_html(el.table_json)
            if table_html:
                parts.append(table_html)
                continue
        text = sanitize(el.text_content or "").strip()
        if not text:
            continue
        if el.element_type == "heading":
            flush_run()
            parts.append(f"<p><b>{html.escape(collapse_hspace(text))}</b></p>")
        elif el.element_type in ("paragraph", "text", "caption"):
            paragraph_run.append(text)
        else:
            flush_run()
            parts.append(f"<p>{html.escape(collapse_hspace(text)).replace(chr(10), '<br>')}</p>")
    flush_run()
    joined = "".join(parts)
    return joined if joined.strip() else None


def _table_html(table_json: dict) -> str | None:
    """Docling table structure → real HTML table (all cells escaped)."""
    try:
        columns = table_json.get("columns") or []
        rows = table_json.get("data") or []
        if not isinstance(columns, list):
            columns = []
        if not isinstance(rows, list) or any(
            not isinstance(r, (list, tuple)) for r in rows
        ):
            rows = []
        if not rows and not columns:
            return None
        head = "".join(f"<th>{html.escape(str(c))}</th>" for c in columns)
        body = "".join(
            "<tr>" + "".join(f"<td>{html.escape(sanitize(str(v)))}</td>" for v in row) + "</tr>"
            for row in rows
        )
        thead = f"<tr>{head}</tr>" if head else ""
        return f"<table>{thead}{body}</table>"
    except Exception:  # noqa: BLE001 — malformed table json degrades to nothing
        return None


def _pptx_pages_html(source_path) -> dict[int, str]:
    from pptx import Presentation

    prs = Presentation(str(source_path))
    # Pass 1: find slide-master boilerplate (footers, "Slide N") by frequency
    per_page_keys: dict[int, set[str]] = {}
    for idx, slide in enumerate(prs.slides, start=1):
        keys: set[str] = set()
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                raw = "".join(run.text for run in para.runs).strip()
                if raw and len(raw) <= REPEAT_MAX_LEN:
                    keys.add(repeat_key(raw))
        per_page_keys[idx] = keys
    skip = repeated_keys(per_page_keys)

    out: dict[int, str] = {}
    for idx, slide in enumerate(prs.slides, start=1):
        parts: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            for para in shape.text_frame.paragraphs:
                raw = "".join(run.text for run in para.runs)
                if _is_repeat(raw, skip):
                    continue
                spans = [
                    _span_html(
                        run.text,
                        bool(run.font.bold),
                        bool(run.font.italic),
                        any(m in (run.font.name or "").lower() for m in _MONO_FONT),
                    )
                    for run in para.runs
                    if run.text.strip()
                ]
                if spans:
                    parts.append(f"<p>{''.join(spans)}</p>")
        if parts:
            out[idx] = "".join(parts)
    return out


def build_text_html(db: Session, document_id: int) -> int:
    """Populate document_pages.text_html for one document. Idempotent."""
    doc = db.get(Document, document_id)
    if doc is None:
        return 0
    pages = (
        db.execute(
            select(DocumentPage)
            .where(DocumentPage.document_id == document_id)
            .order_by(DocumentPage.page_no)
        )
        .scalars()
        .all()
    )
    elements_by_page: dict[int, list[DocElement]] = {}
    for el in (
        db.execute(
            select(DocElement)
            .where(DocElement.document_id == document_id)
            .order_by(DocElement.order_index)
        )
        .scalars()
        .all()
    ):
        elements_by_page.setdefault(el.page_id, []).append(el)

    updated = 0
    if doc.kind == DocumentKind.pptx:
        by_no = _pptx_pages_html(files.resolve(doc.storage_path))
        for page in pages:
            if (page.text_html or "").startswith(USER_EDITED_MARKER):
                continue  # never clobber the user's manual cleanup
            page_elements = elements_by_page.get(page.id, [])
            if any(el.element_type == "table" for el in page_elements):
                html_text = _elements_page_html(page_elements)
            else:
                html_text = by_no.get(page.page_no) or _elements_page_html(page_elements)
            if html_text:
                page.text_html = html_text
                updated += 1
    else:
        import fitz

        from manabi_server.processing.pipeline import parse_source_path

        # Read the normalized (spread-split) PDF so native spans line up with
        # the same logical pages the elements/renders use.
        with fitz.open(parse_source_path(doc)) as pdf:
            skip = repeated_keys(
                {no: _pdf_page_line_keys(pdf[no]) for no in range(pdf.page_count)}
            )
            for page in pages:
                if (page.text_html or "").startswith(USER_EDITED_MARKER):
                    continue
                page_elements = elements_by_page.get(page.id, [])
                # Pages with detected tables use the elements path so
                # table_json renders as a real <table> in reading order.
                if any(el.element_type == "table" for el in page_elements):
                    html_text = _elements_page_html(page_elements)
                else:
                    native = None
                    if 1 <= page.page_no <= pdf.page_count:
                        native = _pdf_page_html(pdf[page.page_no - 1], skip)
                    html_text = native or _elements_page_html(page_elements)
                    if html_text:
                        html_text += _figure_markers(page_elements)
                if html_text:
                    page.text_html = html_text
                    updated += 1
    db.commit()
    return updated
