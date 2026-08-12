"""Document processing pipeline — runs SYNC inside the CPU worker's task
thread. Stages checkpoint via documents.extract_stage so a crashed job
resumes at the failed stage instead of redoing everything. Each stage is
wipe-and-redo idempotent for its own artifacts. The original upload is never
touched.

Stages: structure (Docling parse + speaker notes + persist pages/elements)
      → render (PNGs per page; LibreOffice for PPTX, non-fatal on failure)
      → chunk  (structure-aware chunks + FTS via generated tsv column)
"""

import logging
import re
import subprocess
import tempfile
from pathlib import Path

from manabi_core.models import (
    Chunk,
    DocElement,
    Document,
    DocumentKind,
    DocumentPage,
    ExtractStatus,
    Job,
    JobStatus,
    Module,
)
from sqlalchemy import delete, select
from sqlalchemy.orm import Session

from manabi_server.config import get_settings
from manabi_server.processing import chunking
from manabi_server.storage import files

log = logging.getLogger("manabi.pipeline")

STAGES = ["structure", "render", "chunk", "embed"]


def _progress(db: Session, job: Job | None, pct: int, note: str) -> None:
    if job is not None:
        job.progress_pct = pct
        job.progress_note = note
        db.commit()


def run_pipeline(db: Session, document_id: int, job_id: int | None) -> None:
    doc = db.get(Document, document_id)
    job = db.get(Job, job_id) if job_id else None
    if doc is None or doc.deleted_at is not None:
        if job is not None:
            job.status = JobStatus.cancelled
            job.progress_note = "document deleted"
            db.commit()
        return

    doc.extract_status = ExtractStatus.processing
    doc.error = None
    if job is not None:
        from datetime import UTC, datetime

        job.status = JobStatus.running
        job.started_at = datetime.now(UTC)
    db.commit()

    done = STAGES.index(doc.extract_stage) + 1 if doc.extract_stage in STAGES else 0
    try:
        if doc.processing_mode == "render_only":
            # Store + view only (e.g. syllabus): render pages, skip parse,
            # text extraction, chunking, and embeddings entirely.
            _progress(db, job, 20, "Rendering pages")
            db.execute(delete(Chunk).where(Chunk.document_id == doc.id))
            db.execute(delete(DocElement).where(DocElement.document_id == doc.id))
            db.commit()
            _stage_render(db, doc, job)
            doc.extract_stage = "embed"  # terminal — nothing left to resume
            db.commit()
        if doc.processing_mode != "render_only" and done < 1:
            _progress(db, job, 10, "Parsing document")
            _stage_structure(db, doc)
            doc.extract_stage = "structure"
            db.commit()
        if doc.processing_mode != "render_only" and done < 2:
            _progress(db, job, 45, "Rendering pages")
            _stage_render(db, doc, job)
            from manabi_server.processing.text_html import build_text_html

            build_text_html(db, doc.id)
            doc.extract_stage = "render"
            db.commit()
        if doc.processing_mode != "render_only" and done < 3:
            _progress(db, job, 80, "Creating chunks")
            _stage_chunk(db, doc)
            doc.extract_stage = "chunk"
            db.commit()
        if doc.processing_mode != "render_only" and done < 4:
            _progress(db, job, 92, "Indexing for retrieval")
            from manabi_server.processing.embedding import embed_missing_chunks_sync

            embed_missing_chunks_sync(db)
            doc.extract_stage = "embed"
            db.commit()

        doc.extract_status = ExtractStatus.ready
        module = db.get(Module, doc.module_id)
        if module is not None:
            module.content_version += 1
        if job is not None:
            job.status = JobStatus.succeeded
            job.progress_pct = 100
            job.progress_note = "Ready"
        db.commit()
    except Exception as exc:
        db.rollback()
        log.exception("pipeline failed for document %s", document_id)
        doc.extract_status = ExtractStatus.failed
        doc.error = f"{type(exc).__name__}: {str(exc)[:500]}"
        if job is not None:
            job.status = JobStatus.failed
            job.error = doc.error
        db.commit()
        raise


# ── Stage 1: structure ────────────────────────────────────────────────────


def parse_source_path(doc: Document) -> Path:
    """The file the extraction pipeline reads: the spread-split normalized PDF
    when present, else the true original. The 'Original' view/download must NOT
    use this — it always serves doc.storage_path."""
    if doc.kind == DocumentKind.pdf:
        norm = files.resolve(files.normalized_path(doc.id))
        if norm.exists():
            return norm
    return files.resolve(doc.storage_path)


def _ensure_normalized_pdf(doc: Document) -> str:
    """Write storage/normalized/<id>.pdf when the source PDF needs a normalized
    copy for extraction, and return a cache tag naming the transformation the
    parse actually reads:
      'spread' — split two-page spreads into portrait pages (user opt-in;
                 any page rotation is baked upright as part of the split);
      'rot'    — page rotation baked upright (no split) so parse + render + bbox
                 share one coordinate space and region overlays stay aligned;
      'orig'   — no normalization; extraction reads the untouched upload;
      'na'     — not a PDF.
    Under 'auto' we detect spreads and record `detected_layout` as a UI hint but
    never cut the file automatically. The tag (not `detected_layout`) keys the
    parse cache so an uncut auto-detected spread and an opted-in split — which
    read different source bytes — never collide."""
    if doc.kind != DocumentKind.pdf:
        return "na"
    from manabi_server.processing import layout

    target = files.resolve(files.normalized_path(doc.id))

    def _finish_no_split(src, detected: str) -> str:
        """Bake rotation into the normalized file when any page is rotated, else
        remove it and read the original. Returns the cache tag ('rot'/'orig')."""
        doc.detected_layout = detected
        baked = layout.normalize_rotation(src)
        if baked is None:
            target.unlink(missing_ok=True)
            return "orig"
        try:
            target.parent.mkdir(parents=True, exist_ok=True)
            tmp = target.with_suffix(".pdf.tmp")
            baked.save(str(tmp))
            tmp.replace(target)
        finally:
            baked.close()
        log.info("baked page rotation upright for doc %s", doc.id)
        return "rot"

    import pymupdf

    with pymupdf.open(files.resolve(doc.storage_path)) as src:
        force = doc.page_layout == "spread"
        if doc.page_layout == "single":
            return _finish_no_split(src, "single")
        # Native-text PDFs are never spreads (skip the whole-doc rasterization).
        if not force and layout.has_native_text(src):
            return _finish_no_split(src, "single")
        is_spread_detected, gutters = layout.detect_spread(src)
        if not force:
            # 'auto' → only hint; the user cuts explicitly via the layout control.
            return _finish_no_split(src, "spread" if is_spread_detected else "single")
        # Opted in: guarantee a cut on every page (centre where no clean gutter).
        gutters = [
            g if g is not None else pg.rect.width / 2
            for g, pg in zip(gutters, src, strict=False)
        ]
        doc.detected_layout = "spread"
        out = layout.split_spread_pdf(src, gutters)
        try:
            target.parent.mkdir(parents=True, exist_ok=True)
            tmp = target.with_suffix(".pdf.tmp")
            out.save(str(tmp))
            tmp.replace(target)
        finally:
            out.close()
    log.info("split doc %s into portrait pages (user-selected spread)", doc.id)
    return "spread"


def _stage_structure(db: Session, doc: Document) -> None:
    # wipe-and-redo (also removes dependent chunks via cascade-by-hand)
    db.execute(delete(Chunk).where(Chunk.document_id == doc.id))
    db.execute(delete(DocElement).where(DocElement.document_id == doc.id))
    db.execute(delete(DocumentPage).where(DocumentPage.document_id == doc.id))
    db.commit()

    # Normalize the source first (split spreads and/or bake page rotation
    # upright); every stage below then reads the normalized file so its 1:1
    # page_no and upright coordinate space hold.
    norm_tag = _ensure_normalized_pdf(doc)
    source = parse_source_path(doc)
    # Separator must be filesystem-safe (no ':' — invalid on Windows). The tag
    # names the exact source transformation the parse reads (orig/rot/spread),
    # so different sources never share a cache entry.
    parsed = _docling_parse(source, cache_key=f"{doc.content_hash}-{norm_tag}")
    parsed["elements"] = _fix_reading_order(
        parsed["elements"], aggressive_columns=doc.detected_layout == "spread"
    )
    parsed["elements"] = _strip_boilerplate(parsed["elements"])
    if doc.detected_layout == "spread":
        # Running heads on scanned books OCR differently each page, defeating
        # exact-repeat stripping — remove them by position instead.
        parsed["elements"] = _strip_margin_repeats(parsed["elements"])
    # Re-join scanned drop-cap initials ("S creams" → "Screams"); safe no-op on
    # native text and non-drop-cap paragraphs.
    parsed["elements"] = _join_drop_caps(parsed["elements"])

    notes_by_page: dict[int, str] = {}
    titles_by_page: dict[int, str] = {}
    if doc.kind == DocumentKind.pptx:
        notes_by_page, titles_by_page = _pptx_notes_and_titles(source)
    else:
        # Derive page titles from SURVIVING headings (post-boilerplate strip),
        # so a repeated running head can never become a page title/heading_path.
        for el in parsed["elements"]:
            if el["type"] == "heading" and el["text"]:
                titles_by_page.setdefault(el["page_no"], el["text"][:512])

    page_count = max(
        [p for p in parsed["pages"]] + list(notes_by_page) + list(titles_by_page),
        default=1,
    )
    pages: dict[int, DocumentPage] = {}
    for page_no in range(1, page_count + 1):
        page = DocumentPage(
            document_id=doc.id,
            page_no=page_no,
            title=titles_by_page.get(page_no),
            speaker_notes=notes_by_page.get(page_no),
        )
        db.add(page)
        pages[page_no] = page
    db.flush()

    for order_index, item in enumerate(parsed["elements"]):
        db.add(
            DocElement(
                document_id=doc.id,
                page_id=pages[item["page_no"]].id if item["page_no"] in pages else None,
                order_index=order_index,
                element_type=item["type"],
                text_content=item["text"] or None,
                table_json=item.get("table"),
                bbox=item.get("bbox"),
            )
        )
    doc.page_count = page_count
    db.commit()


def _strip_boilerplate(elements: list[dict]) -> list[dict]:
    """Drop repeating header/footer lines (copyright notices, 'Page N', …)
    that appear on a large fraction of pages. Runs before persist, so the
    cleanup reaches elements, chunks, and AI context alike. build_text_html
    applies the same rules (shared helpers) to the native text-layer views."""
    from manabi_server.processing.text_html import (
        REPEAT_MAX_LEN,
        repeat_key,
        repeated_keys,
    )

    per_page: dict[int, set[str]] = {}
    for el in elements:
        keys = per_page.setdefault(el["page_no"], set())
        text = (el.get("text") or "").strip()
        if text and len(text) <= REPEAT_MAX_LEN:
            keys.add(repeat_key(text))

    boilerplate = repeated_keys(per_page)
    if not boilerplate:
        return elements

    kept = [
        el
        for el in elements
        if not (
            (el.get("text") or "").strip()
            and len((el.get("text") or "").strip()) <= REPEAT_MAX_LEN
            and repeat_key(el["text"]) in boilerplate
        )
    ]
    log.info(
        "stripped %d boilerplate elements (%d distinct lines across %d pages)",
        len(elements) - len(kept), len(boilerplate), len(per_page),
    )
    return kept


MARGIN_REPEAT_TEXT_MAX = 60      # running heads / page numbers are short
MARGIN_REPEAT_PAGE_RATIO = 0.4   # …and appear on a large fraction of pages


def _strip_margin_repeats(elements: list[dict]) -> list[dict]:
    """Drop the top-most and/or bottom-most short element on each page when that
    position repeats across the document — running heads and page numbers on
    scanned books. Position-based, so it catches OCR variants that exact-text
    boilerplate stripping misses. Spread-gated by the caller (never PPTX)."""
    by_page: dict[int, list[dict]] = {}
    for el in elements:
        by_page.setdefault(el["page_no"], []).append(el)
    if len(by_page) < 5:
        return elements

    top_ids: set[int] = set()
    bottom_ids: set[int] = set()
    top_hits = bottom_hits = 0
    for items in by_page.values():
        boxed = [el for el in items if el.get("bbox")]
        if len(boxed) < 3:
            continue
        # bottom-left origin: larger t = higher on the page
        topmost = max(boxed, key=lambda el: el["bbox"]["t"])
        bottommost = min(boxed, key=lambda el: el["bbox"]["t"])
        if len(topmost.get("text") or "") <= MARGIN_REPEAT_TEXT_MAX:
            top_ids.add(id(topmost))
            top_hits += 1
        if len(bottommost.get("text") or "") <= MARGIN_REPEAT_TEXT_MAX:
            bottom_ids.add(id(bottommost))
            bottom_hits += 1

    pages = len(by_page)
    strip: set[int] = set()
    if top_hits / pages >= MARGIN_REPEAT_PAGE_RATIO:
        strip |= top_ids
    if bottom_hits / pages >= MARGIN_REPEAT_PAGE_RATIO:
        strip |= bottom_ids
    if not strip:
        return elements
    kept = [el for el in elements if id(el) not in strip]
    log.info("stripped %d margin running-heads/footers", len(elements) - len(kept))
    return kept


# A scanned drop-cap initial is read by full-page OCR as a lone capital split
# from its word ("S creams"). Re-join it — but ONLY a consonant excluding the
# real single-letter words A/I/O, so "A man", "I have", "O nce" are never
# touched. Position + length are further gated by the caller.
_DROP_CAP_RE = re.compile(r"^([B-HJ-NP-Z]) (?=[a-z])")
_DROP_CAP_MIN_LEN = 40  # drop caps lead real body prose, not short labels


def _join_drop_caps(elements: list[dict]) -> list[dict]:
    """Re-attach a drop-cap initial to its word in the first body paragraph
    right after a heading — the classic drop-cap position. Gated tightly (post-
    heading + long paragraph + consonant≠A/I/O) so ordinary text and technical
    openers ('T cell') are never merged. No-op when OCR already joined it."""
    prev_heading = False
    for el in elements:
        text = el.get("text") or ""
        if (
            el.get("type") == "paragraph"
            and prev_heading
            and len(text) >= _DROP_CAP_MIN_LEN
        ):
            joined = _DROP_CAP_RE.sub(r"\1", text, count=1)
            if joined != text:
                el["text"] = joined
                log.info("re-joined drop-cap initial: %r → %r", text[:14], joined[:12])
        prev_heading = el.get("type") == "heading"
    return elements


INVERSION_TOLERANCE_PTS = 5
# Scrambled OCR pages measure ~20%+; genuine multi-column pages have one
# inversion per column break (a few % with realistic element counts).
INVERSION_RATE_THRESHOLD = 0.15
# A column gap must be at least this fraction of the content width to count.
COLUMN_GAP_MIN_FRAC = 0.12
# Elements wider than this fraction of content width span columns (headings,
# tables, figures) and act as row separators rather than column members.
SPAN_WIDTH_FRAC = 0.6


def _detect_column_split(boxed: list[dict], content_l: float, content_r: float) -> float | None:
    """If a page reads as two columns, return the x that separates them, else
    None. Looks for a wide gap between element centres in the middle of the
    content, with real text on both sides (so figures/tables don't trip it)."""
    width = content_r - content_l
    if width <= 0:
        return None
    centers = sorted(
        (el["bbox"]["l"] + el["bbox"]["r"]) / 2
        for el in boxed
        if (el["bbox"]["r"] - el["bbox"]["l"]) < width * SPAN_WIDTH_FRAC
    )
    if len(centers) < 4:
        return None
    best_gap, best_mid = 0.0, None
    for a, b in zip(centers, centers[1:], strict=False):
        mid = (a + b) / 2
        if content_l + width * 0.2 < mid < content_l + width * 0.8 and (b - a) > best_gap:
            best_gap, best_mid = b - a, mid
    if best_mid is None or best_gap < width * COLUMN_GAP_MIN_FRAC:
        return None
    left = sum(1 for c in centers if c < best_mid)
    right = len(centers) - left
    return best_mid if left >= 2 and right >= 2 else None


def _fix_reading_order(
    elements: list[dict], aggressive_columns: bool = False
) -> list[dict]:
    """Repair scrambled per-page element order from OCR/layout analysis.

    Two-column pages (common in scanned books) are ordered column-by-column,
    top-to-bottom — a plain top-to-bottom sort would interleave the columns.
    But re-ordering a native PDF that docling already got right is risky, so
    column ordering only runs when `aggressive_columns` (spread scans) OR the
    page's element sequence already looks scrambled (high inversion rate).
    Single-column pages keep the same conservative gate."""
    by_page: dict[int, list[dict]] = {}
    for el in elements:
        by_page.setdefault(el["page_no"], []).append(el)

    fixed: list[dict] = []
    for page_no in sorted(by_page):
        items = by_page[page_no]
        boxed = [el for el in items if el.get("bbox")]
        if len(boxed) < 4:
            fixed.extend(items)
            continue

        inversions = sum(
            1
            for prev, nxt in zip(boxed, boxed[1:], strict=False)
            if nxt["bbox"]["t"] > prev["bbox"]["t"] + INVERSION_TOLERANCE_PTS
        )
        pairs = max(1, len(boxed) - 1)
        scrambled = inversions / pairs > INVERSION_RATE_THRESHOLD

        content_l = min(el["bbox"]["l"] for el in boxed)
        content_r = max(el["bbox"]["r"] for el in boxed)
        split = _detect_column_split(boxed, content_l, content_r)

        if split is not None and (aggressive_columns or scrambled):
            # Column-aware: left column fully (top→bottom), then right column.
            # Full-width spanning elements (headings) fall into the left band by
            # their centre, keeping them near their vertical position.
            width = content_r - content_l

            def col_key(el: dict, split: float = split, width: float = width) -> tuple:
                if not el.get("bbox"):
                    return (0, 0.0)
                b = el["bbox"]
                if (b["r"] - b["l"]) >= width * SPAN_WIDTH_FRAC:
                    return (0, -b["t"])  # spanning row stays in the left flow
                center = (b["l"] + b["r"]) / 2
                return (0 if center < split else 1, -b["t"])

            items = sorted(items, key=col_key)
            log.info("ordered page %s as two columns (split x≈%.0f)", page_no, split)
        elif scrambled:
            items = sorted(
                items,
                key=lambda el: (
                    -(el["bbox"]["t"] if el.get("bbox") else 0),
                    el["bbox"]["l"] if el.get("bbox") else 0,
                ),
            )
            log.info(
                "re-sorted scrambled page %s (%d/%d inversions)",
                page_no, inversions, pairs,
            )
        fixed.extend(items)
    return fixed


_converters: dict[bool, object] = {}


def _get_converter(full_page_ocr: bool = False):
    """Docling converter singletons (one per OCR mode) — model initialization is
    expensive (tens of seconds); pay it once per mode per worker process.

    full_page_ocr re-OCRs the whole page image instead of only the regions the
    layout model flags as needing it. It recovers large decorative initials
    (drop caps / chapter caps) that region OCR drops, so scanned books read
    correctly ("ONE", "Screams" rather than "NE", "creams"). It's slower and
    would override a real text layer, so it's used ONLY for scans — native-text
    PDFs keep the fast, exact default path."""
    if full_page_ocr not in _converters:
        import os

        from docling.datamodel.base_models import InputFormat
        from docling.datamodel.pipeline_options import (
            AcceleratorOptions,
            PdfPipelineOptions,
        )
        from docling.document_converter import DocumentConverter, PdfFormatOption

        # Docling defaults to 4 threads; the layout/OCR/table models are the
        # pipeline's dominant cost, so use the whole machine.
        opts = PdfPipelineOptions(
            accelerator_options=AcceleratorOptions(num_threads=os.cpu_count() or 4)
        )
        if full_page_ocr:
            try:
                from docling.datamodel.pipeline_options import OcrMode

                opts.ocr_options.mode = OcrMode.FULL_PAGE
            except Exception:  # noqa: BLE001 — older docling: deprecated flag
                opts.ocr_options.force_full_page_ocr = True
        _converters[full_page_ocr] = DocumentConverter(
            format_options={InputFormat.PDF: PdfFormatOption(pipeline_options=opts)}
        )
    return _converters[full_page_ocr]


# Bump whenever parse or normalization logic changes (e.g. rotation baking,
# reading-order, boilerplate rules) so stale cached parses are auto-invalidated
# instead of silently re-served on the next re-extraction.
#   v2: rotated scans rasterized upright in normalize_rotation (was sideways).
#   v3: scanned PDFs use full-page OCR (recovers dropped drop-cap initials).
_PARSE_CACHE_VERSION = 3


def _parse_cache_path(cache_key: str) -> Path:
    return (
        files.storage_root()
        / "parse-cache"
        / f"v{_PARSE_CACHE_VERSION}-{cache_key}.json"
    )


def _docling_parse(source: Path, cache_key: str | None = None) -> dict:
    """Parse via Docling, memoized on the file's content hash: the model pass
    is deterministic per file, so retries and re-extractions (e.g. after a
    chunking or cleanup change) skip minutes of CPU work."""
    import json

    cache_path = _parse_cache_path(cache_key) if cache_key else None
    if cache_path is not None and cache_path.exists():
        try:
            raw = json.loads(cache_path.read_text(encoding="utf-8"))
            return {
                "elements": raw["elements"],
                "titles": {int(k): v for k, v in raw["titles"].items()},
                "pages": set(raw["pages"]),
            }
        except Exception:  # noqa: BLE001 — corrupt cache → re-parse
            log.warning("parse cache unreadable, re-parsing: %s", cache_path)

    # Scans (no usable text layer) get full-page OCR so drop-cap initials survive;
    # native-text PDFs use the fast default path (its real text layer is exact).
    import pymupdf

    from manabi_server.processing import layout

    with pymupdf.open(str(source)) as _probe:
        is_scan = not layout.has_native_text(_probe)
    result = _get_converter(full_page_ocr=is_scan).convert(str(source))
    dl_doc = result.document

    elements: list[dict] = []
    titles: dict[int, str] = {}
    pages_seen: set[int] = set()

    for item, _level in dl_doc.iterate_items():
        label = str(getattr(item, "label", "") or "").lower()
        page_no = 1
        bbox = None
        prov = getattr(item, "prov", None)
        if prov:
            page_no = getattr(prov[0], "page_no", 1) or 1
            bb = getattr(prov[0], "bbox", None)
            if bb is not None:
                bbox = {"l": bb.l, "t": bb.t, "r": bb.r, "b": bb.b}
        pages_seen.add(page_no)

        from manabi_server.processing.text_html import sanitize

        text = sanitize(getattr(item, "text", "") or "").strip()
        table = None
        if "table" in label:
            element_type = "table"
            try:
                table = item.export_to_dataframe(dl_doc).to_dict(orient="split")
                text = text or _linearize_table(table)
            except Exception:  # noqa: BLE001 — table extraction is best-effort
                pass
        elif "picture" in label or "figure" in label:
            element_type = "figure"
        elif "section_header" in label or "title" in label or "heading" in label:
            # Demote pseudo-headings (lead-in lines Docling misclassifies,
            # e.g. "can be defined like this:") back to paragraphs.
            if text.endswith(":") or (text and text[0].islower()):
                element_type = "paragraph"
            else:
                element_type = "heading"
                if page_no not in titles and text:
                    titles[page_no] = text[:512]
        elif "list" in label:
            element_type = "list"
        elif "caption" in label:
            element_type = "caption"
        elif "formula" in label:
            element_type = "formula"
        else:
            element_type = "paragraph"

        if not text and element_type not in ("figure",):
            continue
        elements.append(
            {"type": element_type, "text": text, "page_no": page_no, "bbox": bbox, "table": table}
        )

    if cache_path is not None:
        try:
            cache_path.parent.mkdir(parents=True, exist_ok=True)
            cache_path.write_text(
                json.dumps(
                    {
                        "elements": elements,
                        "titles": titles,
                        "pages": sorted(pages_seen),
                    }
                ),
                encoding="utf-8",
            )
        except Exception:  # noqa: BLE001 — cache is an optimization only
            log.warning("could not write parse cache: %s", cache_path)

    return {"elements": elements, "titles": titles, "pages": pages_seen}


def _linearize_table(table: dict) -> str:
    try:
        header = " | ".join(str(c) for c in table.get("columns", []))
        rows = "\n".join(" | ".join(str(v) for v in row) for row in table.get("data", []))
        return f"{header}\n{rows}".strip()
    except Exception:  # noqa: BLE001
        return ""


def _pptx_notes_and_titles(source: Path) -> tuple[dict[int, str], dict[int, str]]:
    from pptx import Presentation

    from manabi_server.processing.text_html import sanitize

    notes: dict[int, str] = {}
    titles: dict[int, str] = {}
    prs = Presentation(str(source))
    for idx, slide in enumerate(prs.slides, start=1):
        try:
            if slide.shapes.title is not None and slide.shapes.title.text.strip():
                titles[idx] = sanitize(slide.shapes.title.text).strip()[:512]
        except Exception:  # noqa: BLE001
            pass
        if slide.has_notes_slide:
            text = sanitize(slide.notes_slide.notes_text_frame.text).strip()
            if text:
                notes[idx] = text
    return notes, titles


# ── Stage 2: render ───────────────────────────────────────────────────────


def _stage_render(db: Session, doc: Document, job: Job | None) -> None:
    source = parse_source_path(doc)  # normalized (split) PDF if present
    pdf_path = source

    if doc.kind == DocumentKind.pptx:
        converted = _pptx_to_pdf(source)
        if converted is None:
            log.warning("LibreOffice conversion failed for doc %s — text-only viewer", doc.id)
            return  # non-fatal: pages keep render_path NULL
        pdf_path = converted

    # Drop renders from a previous extraction (a re-split can change the page
    # count, e.g. spread→single), so no orphaned high-numbered PNGs linger.
    import shutil

    import fitz  # PyMuPDF

    for sub in (f"renders/{doc.id}", f"thumbs/{doc.id}"):
        shutil.rmtree(files.resolve(sub), ignore_errors=True)

    pages = (
        db.execute(
            select(DocumentPage)
            .where(DocumentPage.document_id == doc.id)
            .order_by(DocumentPage.page_no)
        )
        .scalars()
        .all()
    )
    by_no = {p.page_no: p for p in pages}

    with fitz.open(pdf_path) as pdf:
        total = pdf.page_count
        for i, pdf_page in enumerate(pdf):
            page_no = i + 1
            row = by_no.get(page_no)
            if row is None:
                row = DocumentPage(document_id=doc.id, page_no=page_no)
                db.add(row)
                db.flush()
                by_no[page_no] = row
            pix = pdf_page.get_pixmap(matrix=fitz.Matrix(2, 2))
            rel = files.render_path(doc.id, page_no)
            files.write_atomic(rel, pix.tobytes("png"))
            thumb = pdf_page.get_pixmap(matrix=fitz.Matrix(0.35, 0.35))
            rel_thumb = files.thumb_path(doc.id, page_no)
            files.write_atomic(rel_thumb, thumb.tobytes("png"))
            row.render_path = rel
            row.thumb_path = rel_thumb
            row.width = pix.width
            row.height = pix.height
            if page_no % 5 == 0 or page_no == total:
                pct = 45 + int(40 * page_no / total)
                _progress(db, job, pct, f"Rendering pages {page_no}/{total}")
        doc.page_count = total
    if pdf_path != source:
        pdf_path.unlink(missing_ok=True)
    db.commit()


def _pptx_to_pdf(source: Path) -> Path | None:
    settings = get_settings()
    outdir = Path(tempfile.mkdtemp(prefix="manabi_soffice_"))
    try:
        subprocess.run(
            [
                settings.soffice_path,
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(outdir),
                str(source),
            ],
            check=True,
            capture_output=True,
            timeout=300,
        )
        produced = list(outdir.glob("*.pdf"))
        return produced[0] if produced else None
    except Exception:  # noqa: BLE001 — render fallback handles it
        log.exception("soffice conversion failed")
        return None


# ── Stage 3: chunk ────────────────────────────────────────────────────────


def _stage_chunk(db: Session, doc: Document) -> None:
    db.execute(delete(Chunk).where(Chunk.document_id == doc.id))

    page_rows = (
        db.execute(select(DocumentPage).where(DocumentPage.document_id == doc.id))
        .scalars()
        .all()
    )
    element_rows = (
        db.execute(
            select(DocElement)
            .where(DocElement.document_id == doc.id)
            .order_by(DocElement.order_index)
        )
        .scalars()
        .all()
    )
    page_no_by_id = {p.id: p.page_no for p in page_rows}
    elements = [
        chunking.ElementIn(
            id=e.id,
            page_no=page_no_by_id.get(e.page_id, 1),
            element_type=e.element_type,
            text=e.text_content or "",
        )
        for e in element_rows
    ]

    if doc.kind == DocumentKind.pptx:
        pages = [
            chunking.PageIn(page_no=p.page_no, title=p.title, speaker_notes=p.speaker_notes)
            for p in page_rows
        ]
        results = chunking.chunk_pptx(pages, elements)
    else:
        results = chunking.chunk_pdf(elements)

    for c in results:
        db.add(
            Chunk(
                module_id=doc.module_id,
                document_id=doc.id,
                page_start=c.page_start,
                page_end=c.page_end,
                element_ids=c.element_ids,
                heading_path=c.heading_path,
                text=c.text,
                token_count=c.token_count,
                content_hash=c.content_hash,
            )
        )
    db.commit()
